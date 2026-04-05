import streamlit as st
import cv2
import mediapipe as mp
import pandas as pd
import numpy as np
import time
import av
import io
from PIL import Image
from pptx import Presentation
from streamlit_webrtc import webrtc_streamer, VideoProcessorBase, RTCConfiguration, WebRtcMode

# --- Configuration ---
st.set_page_config(
    page_title="GestureAI - Hand Gesture Control",
    page_icon="⚡",
    layout="wide",
    initial_sidebar_state="expanded",
)

# --- Custom CSS ---
st.markdown("""
    <style>
    .main { background-color: #ffffff; }
    .stButton>button { width: 100%; border-radius: 12px; height: 3em; font-weight: bold; }
    .stSidebar { background-color: #f8f9fa; }
    </style>
    """, unsafe_allow_html=True)

# --- MediaPipe Setup ---
@st.cache_resource
def get_mediapipe_hands():
    return mp.solutions.hands.Hands(
        static_image_mode=False,
        max_num_hands=1,
        min_detection_confidence=0.7,
        min_tracking_confidence=0.5
    )

mp_hands = mp.solutions.hands
mp_draw = mp.solutions.drawing_utils

import queue

# --- Heuristic Gesture Logic ---
def calculate_distance(p1, p2):
    return np.sqrt((p1.x - p2.x)**2 + (p1.y - p2.y)**2 + (p1.z - p2.z)**2)

def detect_gesture_heuristic(hand_landmarks, prev_x=None):
    thumb_tip = hand_landmarks.landmark[4]
    index_tip = hand_landmarks.landmark[8]
    middle_tip = hand_landmarks.landmark[12]
    index_mcp = hand_landmarks.landmark[5]
    middle_mcp = hand_landmarks.landmark[9]
    
    # Tọa độ X trung tâm bàn tay (dùng landmark 9 - Middle Finger MCP)
    curr_x = hand_landmarks.landmark[9].x
    
    dist_thumb_index = calculate_distance(thumb_tip, index_tip)
    is_index_ext = index_tip.y < index_mcp.y
    is_middle_folded = middle_tip.y > middle_mcp.y
    
    gesture = "None"
    if prev_x is not None:
        diff = curr_x - prev_x
        if diff > 0.08:      # === FIX: giảm ngưỡng ===
            gesture = "Swipe Right"
        elif diff < -0.08:   # === FIX: giảm ngưỡng ===
            gesture = "Swipe Left"
    
    # ... phần còn lại giữ nguyên
    return gesture, curr_x
    
    if gesture == "None":
        if dist_thumb_index < 0.05:
            gesture = "Click"
        elif is_index_ext and is_middle_folded:
            gesture = "Laser Pointer"
        elif not is_index_ext and is_middle_folded:
            gesture = "Fist (Nắm tay)"
        else:
            gesture = "Open Hand"
            
    return gesture, curr_x

# --- WebRTC Processor ---
    class VideoProcessor(VideoProcessorBase):
        def __init__(self):
            self.hands = get_mediapipe_hands()
            self.result_queue = queue.Queue()
            self.prev_x = None
            self.last_gesture_time = 0
            self.last_gesture = "None"          # ← thêm dòng này
    
        def recv(self, frame):
            img = frame.to_ndarray(format="bgr24")
            img = cv2.flip(img, 1)
            rgb_img = cv2.cvtColor(img, cv2.COLOR_BGR2RGB)
            results = self.hands.process(rgb_img)
    
            if results.multi_hand_landmarks:
                for hand_landmarks in results.multi_hand_landmarks:
                    mp_draw.draw_landmarks(img, hand_landmarks, mp_hands.HAND_CONNECTIONS)
                    gesture, curr_x = detect_gesture_heuristic(hand_landmarks, self.prev_x)
                    self.prev_x = curr_x
                    self.last_gesture = gesture                # ← thêm dòng này
                    
                    color = (0, 255, 0)
                    if "Swipe" in gesture:
                        color = (0, 165, 255)
                    cv2.putText(img, gesture, (50, 50), cv2.FONT_HERSHEY_SIMPLEX, 1, color, 2)
    
                    curr_time = time.time()
                    if ("Swipe" in gesture) and (curr_time - self.last_gesture_time > 1.0):
                        self.result_queue.put(gesture)
                        self.last_gesture_time = curr_time
            else:
                self.prev_x = None
                self.last_gesture = "None"
            
            return av.VideoFrame.from_ndarray(img, format="bgr24")

# --- Helper for Slide Content ---
    with tab_present:
        if 'prs' not in st.session_state:
            st.warning("Vui lòng tải file PPTX ở tab 'Thiết lập' trước.")
        else:
            if 'slide_idx' not in st.session_state:
                st.session_state.slide_idx = 0
            if 'last_swipe' not in st.session_state:
                st.session_state.last_swipe = None
                st.session_state.swipe_timestamp = 0
    
            st.markdown("### 📺 Đang trình chiếu...")
    
            col_cam, col_slide = st.columns([1, 3])
    
            with col_cam:
                st.write("📷 Camera Control")
                webrtc_ctx = webrtc_streamer(
                    key="present-gesture",
                    mode=WebRtcMode.SENDRECV,
                    rtc_configuration=RTCConfiguration({"iceServers": [{"urls": ["stun:stun.l.google.com:19302"]}]}),
                    media_stream_constraints={"video": {"width": 320, "height": 240}, "audio": False},
                    video_processor_factory=VideoProcessor,
                    async_processing=True,
                )
                st.caption("Vuốt TRÁI ← để về trước | Vuốt PHẢI → để sang sau")
    
                processor = getattr(webrtc_ctx, 'video_processor', None)
                if processor is not None:
                    st.caption(f"**Gesture hiện tại:** {processor.last_gesture}")
    
                    # === FIX CHUYỂN SLIDE (drain hết queue) ===
                    try:
                        while True:                                   # ← quan trọng
                            gesture = processor.result_queue.get_nowait()
                            if "Swipe Right" in gesture:
                                st.session_state.slide_idx = min(len(st.session_state.prs.slides) - 1, st.session_state.slide_idx + 1)
                                st.session_state.last_swipe = "Swipe Right"
                                st.session_state.swipe_timestamp = time.time()
                                st.rerun()
                            elif "Swipe Left" in gesture:
                                st.session_state.slide_idx = max(0, st.session_state.slide_idx - 1)
                                st.session_state.last_swipe = "Swipe Left"
                                st.session_state.swipe_timestamp = time.time()
                                st.rerun()
                    except queue.Empty:
                        pass
    
                st.divider()
                c1, c2 = st.columns(2)
                if c1.button("⬅️ Trước", use_container_width=True):
                    st.session_state.slide_idx = max(0, st.session_state.slide_idx - 1)
                    st.rerun()
                if c2.button("Sau ➡️", use_container_width=True):
                    st.session_state.slide_idx = min(len(st.session_state.prs.slides) - 1, st.session_state.slide_idx + 1)
                    st.rerun()
    
                if st.button("Reset Slide", type="secondary", use_container_width=True):
                    st.session_state.slide_idx = 0
                    st.rerun()
    
            # ==================== SLIDE HIỂN THỊ (đầy đủ + đẹp) ====================
            with col_slide:
                total_slides = len(st.session_state.prs.slides)
                current_idx = st.session_state.slide_idx
                current_slide = st.session_state.prs.slides[current_idx]
    
                st.markdown(f"#### Slide **{current_idx + 1} / {total_slides}**")
    
                # Thông báo swipe giữ 4 giây
                current_time = time.time()
                if st.session_state.last_swipe and current_time - st.session_state.swipe_timestamp < 4.0:
                    icon = "➡️" if "Right" in st.session_state.last_swipe else "⬅️"
                    st.success(f"{icon} **{st.session_state.last_swipe.upper()}** - Slide đã chuyển!", icon=icon)
    
                slide_text = get_slide_content(current_slide)
                
                st.markdown(f"""
                <div style="
                    background: linear-gradient(180deg, #f8f9fa 0%, #ffffff 100%);
                    padding: 50px 40px;
                    border-radius: 20px;
                    border: 4px solid #1f2937;
                    min-height: 560px;
                    max-height: 650px;
                    overflow-y: auto;
                    box-shadow: 0 10px 30px rgba(0,0,0,0.1);
                    color: #1f2937;
                    font-size: 1.45em;
                    line-height: 1.75;
                ">
                    {slide_text}
                </div>
                """, unsafe_allow_html=True)
    
                st.progress((current_idx + 1) / total_slides)
elif page == "Triển khai mô hình":
    st.title("🚀 Triển khai & Trình chiếu")
    
    # Khởi tạo trạng thái slide
    if 'slide_idx' not in st.session_state:
        st.session_state.slide_idx = 0
    if 'last_swipe' not in st.session_state:
        st.session_state.last_swipe = None
        st.session_state.swipe_timestamp = 0

    # Tabs cho các chế độ
    tab_setup, tab_present = st.tabs(["🛠️ Thiết lập", "📺 Chế độ trình chiếu"])

    with tab_setup:
        col_l, col_r = st.columns([1, 1.5])
        with col_l:
            st.subheader("📹 Phân tích Video")
            up_vid = st.file_uploader("Tải video (.mp4, .mov, .avi)", type=["mp4", "mov", "avi"])
            if up_vid:
                st.video(up_vid)
                if st.button("Phân tích Video", type="primary"):
                    st.success("Phân tích hoàn tất!")

        with col_r:
            st.subheader("📂 Quản lý Slide")
            ppt_file = st.file_uploader("Tải file PPTX", type=["pptx"])
            if ppt_file:
                prs_bytes = io.BytesIO(ppt_file.read())
                st.session_state.prs = Presentation(prs_bytes)
                st.success(f"✅ Đã tải: {ppt_file.name} ({len(st.session_state.prs.slides)} slides)")
                st.info("Chuyển sang tab 'Chế độ trình chiếu' để bắt đầu.")

    with tab_present:
        if 'prs' not in st.session_state:
            st.warning("Vui lòng tải file PPTX ở tab 'Thiết lập' trước.")
        else:
            st.markdown("### 📺 Đang trình chiếu...")

            col_cam, col_slide = st.columns([1, 3])

            with col_cam:
                st.write("📷 Camera Control")
                webrtc_ctx = webrtc_streamer(
                    key="present-gesture",
                    mode=WebRtcMode.SENDRECV,
                    rtc_configuration=RTCConfiguration({"iceServers": [{"urls": ["stun:stun.l.google.com:19302"]}]}),
                    media_stream_constraints={"video": {"width": 320, "height": 240}, "audio": False},
                    video_processor_factory=VideoProcessor,
                    async_processing=True,
                )
                st.caption("Vuốt TRÁI ← để về trước | Vuốt PHẢI → để sang sau")

                processor = getattr(webrtc_ctx, 'video_processor', None)
                if processor is not None:
                    st.caption(f"**Gesture hiện tại:** {processor.last_gesture}")

                    # Xử lý swipe (đã fix drain queue)
                    try:
                        while True:
                            gesture = processor.result_queue.get_nowait()
                            if "Swipe Right" in gesture:
                                st.session_state.slide_idx = min(len(st.session_state.prs.slides) - 1, st.session_state.slide_idx + 1)
                                st.session_state.last_swipe = "Swipe Right"
                                st.session_state.swipe_timestamp = time.time()
                                st.rerun()
                            elif "Swipe Left" in gesture:
                                st.session_state.slide_idx = max(0, st.session_state.slide_idx - 1)
                                st.session_state.last_swipe = "Swipe Left"
                                st.session_state.swipe_timestamp = time.time()
                                st.rerun()
                    except queue.Empty:
                        pass

                st.divider()
                c1, c2 = st.columns(2)
                if c1.button("⬅️ Trước", use_container_width=True):
                    st.session_state.slide_idx = max(0, st.session_state.slide_idx - 1)
                    st.rerun()
                if c2.button("Sau ➡️", use_container_width=True):
                    st.session_state.slide_idx = min(len(st.session_state.prs.slides) - 1, st.session_state.slide_idx + 1)
                    st.rerun()

                if st.button("Reset Slide", type="secondary", use_container_width=True):
                    st.session_state.slide_idx = 0
                    st.rerun()

            # ==================== PHẦN HIỂN THỊ SLIDE ====================
            with col_slide:
                total_slides = len(st.session_state.prs.slides)
                current_idx = st.session_state.slide_idx
                current_slide = st.session_state.prs.slides[current_idx]

                st.markdown(f"#### Slide **{current_idx + 1} / {total_slides}**")

                # Thông báo swipe
                current_time = time.time()
                if st.session_state.last_swipe and current_time - st.session_state.swipe_timestamp < 4.0:
                    icon = "➡️" if "Right" in st.session_state.last_swipe else "⬅️"
                    st.success(f"{icon} **{st.session_state.last_swipe.upper()}** - Slide đã chuyển!", icon=icon)

                slide_text = get_slide_content(current_slide)
                
                st.markdown(f"""
                <div style="
                    background: linear-gradient(180deg, #f8f9fa 0%, #ffffff 100%);
                    padding: 50px 40px;
                    border-radius: 20px;
                    border: 4px solid #1f2937;
                    min-height: 560px;
                    max-height: 650px;
                    overflow-y: auto;
                    box-shadow: 0 10px 30px rgba(0,0,0,0.1);
                    color: #1f2937;
                    font-size: 1.45em;
                    line-height: 1.75;
                ">
                    {slide_text}
                </div>
                """, unsafe_allow_html=True)

                st.progress((current_idx + 1) / total_slides)
